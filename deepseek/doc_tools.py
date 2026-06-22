# DeepSeek CLI v7.5 — Document Tools Module
# REAL document processing — NO SIMULATION
# PPTX: create, read, edit, info, slide operations
# XLSX: create, read, edit, info, sheet operations, charts, formulas
# DOCX EDIT: append, replace text, add tables, modify properties
# CSV: read, create, convert
# Conversion: docx→pdf, xlsx→csv, pptx→pdf, etc.
#
# Dependencies (graceful fallback):
#   python-pptx  — PowerPoint files
#   openpyxl     — Excel files
#   python-docx  — Word files (edit)
#   Pillow       — Image handling in PPTX

import os
import json
import re
import io
import csv
import zipfile
import traceback
from pathlib import Path
from datetime import datetime, timedelta
from typing import Any

# ═══════════════════════════════════════════════════════════════
# PPTX STUBS — module-level Emu/Pt for cross-module imports
# ═══════════════════════════════════════════════════════════════
try:
    from pptx.util import Emu, Pt
except ImportError:
    class Emu:
        def __init__(self, value):
            self._value = int(value)
        @property
        def inches(self):
            return self._value / 914400
    def Pt(value):
        return int(value)



# ═══════════════════════════════════════════════════════════════
# PPTX TOOLS — PowerPoint Presentations (REAL)
# ═══════════════════════════════════════════════════════════════

def _check_pptx():
    """Check if python-pptx is available."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        return True
    except ImportError:
        return False


def read_pptx(path: str) -> str:
    """
    Read a PowerPoint (.pptx) file and extract all content.
    Returns slide-by-slide text, notes, shapes, and images info.
    """
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"

    try:
        prs = Presentation(str(p))
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        width_in = Emu(slide_width).inches
        height_in = Emu(slide_height).inches

        lines = []
        lines.append(f"PPTX: {p.name}")
        lines.append(f"File size: {p.stat().st_size / 1024:.1f} KB")
        lines.append(f"Slides: {len(prs.slides)}")
        lines.append(f"Slide size: {width_in:.1f}\" x {height_in:.1f}\" ({prs.slide_width}x{prs.slide_height} EMU)")
        lines.append("=" * 60)

        # Check for slide layouts used
        layout_names = set()
        for slide_layout in prs.slide_layouts:
            layout_names.add(slide_layout.name)

        if layout_names:
            lines.append(f"Available layouts: {', '.join(layout_names)}")
            lines.append("")

        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"\n--- Slide {i} ---")
            if slide.slide_layout:
                lines.append(f"Layout: {slide.slide_layout.name}")

            # Extract all text from shapes
            has_content = False
            for shape in slide.shapes:
                shape_info = _extract_shape_info(shape)
                if shape_info:
                    has_content = True
                    lines.append(shape_info)

            if not has_content:
                lines.append("  (empty slide)")

            # Extract notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append(f"  [Speaker Notes]: {notes_text[:500]}")

        return '\n'.join(lines)

    except Exception as e:
        return f"PPTX read error: {e}\n{traceback.format_exc()}"


def _extract_shape_info(shape) -> str:
    """Extract information from a single shape in a slide."""
    parts = []
    shape_type = shape.shape_type

    # Skip group shapes (too complex to extract simply)
    try:
        stype_name = str(shape_type).replace('MSO_SHAPE_TYPE.', '')
    except Exception:
        stype_name = str(shape_type)

    # Shape name and position
    parts.append(f"  [Shape: {shape.name} | Type: {stype_name}]")

    # Position
    try:
        left_in = Emu(shape.left).inches
        top_in = Emu(shape.top).inches
        width_in = Emu(shape.width).inches
        height_in = Emu(shape.height).inches
        parts.append(f"  Position: left={left_in:.1f}\" top={top_in:.1f}\" size={width_in:.1f}\"x{height_in:.1f}\"")
    except Exception:
        pass

    # Text content
    if shape.has_text_frame:
        text_frame = shape.text_frame
        full_text = text_frame.text.strip()
        if full_text:
            parts.append(f"  Text: {full_text[:1000]}")

            # Paragraph details
            for pi, para in enumerate(text_frame.paragraphs):
                para_text = para.text.strip()
                if para_text and para_text != full_text:
                    alignment = str(para.alignment) if para.alignment else 'default'
                    parts.append(f"    P{pi}: [{alignment}] {para_text[:300]}")

                    # Run-level formatting
                    for ri, run in enumerate(para.runs):
                        run_text = run.text.strip()
                        if run_text and len(run_text) > 2:
                            fmt_parts = []
                            if run.font.bold:
                                fmt_parts.append("bold")
                            if run.font.italic:
                                fmt_parts.append("italic")
                            if run.font.underline:
                                fmt_parts.append("underline")
                            if run.font.size:
                                fmt_parts.append(f"{Pt(run.font.size.pt if hasattr(run.font.size, 'pt') else 0)}pt")
                            if run.font.color and run.font.color.rgb:
                                fmt_parts.append(f"#{run.font.color.rgb}")
                            if fmt_parts:
                                parts.append(f"      R{ri}: [{', '.join(fmt_parts)}] {run_text[:200]}")

    # Table content
    if shape.has_table:
        table = shape.table
        parts.append(f"  Table: {table.rows.__len__()} rows x {len(table.columns)} cols")
        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            if any(cells):
                parts.append(f"    Row {row_idx}: {' | '.join(cells[:15])}")

    # Image info
    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
        try:
            parts.append(f"  Image: content_type={shape.image.content_type}")
        except Exception:
            parts.append("  [Image present]")

    return '\n'.join(parts)


def create_pptx(path: str, title: str = '', slides_json: str = '') -> str:
    """
    Create a PowerPoint presentation from structured data.
    REAL creation with python-pptx — supports layouts, formatting, tables, shapes.

    slides_json format:
    [
      {
        "layout": "title_slide" | "title_content" | "section_header" | "blank",
        "title": "Slide Title",
        "subtitle": "Subtitle text",
        "content": "Body text with \\n for newlines",
        "bullets": ["Point 1", "Point 2", "Point 3"],
        "notes": "Speaker notes",
        "table": {
          "headers": ["Col1", "Col2"],
          "rows": [["val1", "val2"]]
        },
        "background_color": "#FF5733",
        "title_color": "#FFFFFF",
        "title_size": 44,
        "content_size": 18,
        "shapes": [
          {"type": "rect", "text": "Hello", "x": 1, "y": 2, "w": 3, "h": 1,
           "fill_color": "#FF0000", "text_color": "#FFFFFF", "font_size": 14}
        ]
      }
    ]
    """
    out_file = Path(path).expanduser()
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"

    try:
        out_file.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()

        # Parse slides data
        slides = []
        if slides_json:
            try:
                slides = json.loads(slides_json)
            except json.JSONDecodeError as e:
                return f"Error: Invalid JSON in slides: {e}"
        elif title:
            # Simple mode: just title slide + content slides
            slides = [{"layout": "title_slide", "title": title, "subtitle": ""}]

        if not slides:
            # Create minimal presentation
            slide_layout = prs.slide_layouts[6]  # blank
            prs.slides.add_slide(slide_layout)
            prs.save(str(out_file))
            return f"PPTX created (blank): {path} ({out_file.stat().st_size / 1024:.1f} KB)"

        for slide_data in slides:
            layout_name = slide_data.get('layout', 'title_content')
            slide_title = slide_data.get('title', '')
            subtitle = slide_data.get('subtitle', '')
            content = slide_data.get('content', '')
            bullets = slide_data.get('bullets', [])
            notes = slide_data.get('notes', '')
            table_data = slide_data.get('table', None)
            bg_color = slide_data.get('background_color', '')
            title_color = slide_data.get('title_color', '')
            title_size = slide_data.get('title_size', 0)
            content_size = slide_data.get('content_size', 0)
            shapes = slide_data.get('shapes', [])

            # Select layout
            if layout_name == 'title_slide':
                slide_layout = prs.slide_layouts[0]  # Title Slide
            elif layout_name == 'section_header':
                slide_layout = prs.slide_layouts[2]  # Section Header
            elif layout_name == 'blank':
                slide_layout = prs.slide_layouts[6]  # Blank
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content

            slide = prs.slides.add_slide(slide_layout)

            # Set background color
            if bg_color:
                try:
                    bg = slide.background
                    fill = bg.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(bg_color.lstrip('#'))
                except Exception:
                    pass

            # Set title
            if slide_title:
                title_shape = slide.shapes.title
                if title_shape and title_shape.has_text_frame:
                    tf = title_shape.text_frame
                    tf.clear()
                    p_run = tf.paragraphs[0].add_run()
                    p_run.text = slide_title
                    if title_size > 0:
                        p_run.font.size = Pt(title_size)
                    if title_color:
                        try:
                            p_run.font.color.rgb = RGBColor.from_string(title_color.lstrip('#'))
                        except Exception:
                            pass

            # Set subtitle
            if subtitle:
                try:
                    # For title slides, subtitle is in shapes[1]
                    if len(slide.placeholders) > 1:
                        sub_shape = slide.placeholders[1]
                        if sub_shape.has_text_frame:
                            sub_shape.text_frame.paragraphs[0].text = subtitle
                except Exception:
                    pass

            # Add content (body text)
            if content:
                try:
                    body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                    if body_shape and body_shape.has_text_frame:
                        tf = body_shape.text_frame
                        tf.clear()
                        for line in content.split('\n'):
                            para = tf.add_paragraph()
                            run = para.add_run()
                            run.text = line.strip()
                            if content_size > 0:
                                run.font.size = Pt(content_size)
                except Exception:
                    pass

            # Add bullets
            if bullets and isinstance(bullets, list):
                try:
                    body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                    if body_shape and body_shape.has_text_frame:
                        tf = body_shape.text_frame
                        # Clear existing and add bullet points
                        tf.clear()
                        for bullet in bullets:
                            para = tf.add_paragraph()
                            para.level = 0
                            run = para.add_run()
                            run.text = str(bullet)
                            if content_size > 0:
                                run.font.size = Pt(content_size)
                except Exception:
                    # Fallback: add text box with bullets
                    left = Inches(0.5)
                    top = Inches(2.0)
                    width = Inches(9.0)
                    height = Inches(5.0)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    for bullet in bullets:
                        para = tf.add_paragraph()
                        run = para.add_run()
                        run.text = f"  - {bullet}"

            # Add table
            if table_data and isinstance(table_data, dict):
                headers = table_data.get('headers', [])
                rows = table_data.get('rows', [])
                if headers:
                    try:
                        num_rows = len(rows) + 1
                        num_cols = len(headers)
                        left = Inches(0.5)
                        top = Inches(3.0) if bullets else Inches(2.0)
                        width = Inches(9.0)
                        height = Inches(0.5 * num_rows)

                        table_shape = slide.shapes.add_table(
                            num_rows, num_cols, left, top, width, height
                        )
                        table = table_shape.table

                        # Header row
                        for j, h in enumerate(headers):
                            cell = table.cell(0, j)
                            cell.text = str(h)
                            # Style header
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.bold = True
                                paragraph.font.size = Pt(12)

                        # Data rows
                        for i, row in enumerate(rows):
                            for j, cell_val in enumerate(row):
                                if j < num_cols:
                                    table.cell(i + 1, j).text = str(cell_val)
                    except Exception as e:
                        pass  # Skip table on error

            # Add custom shapes
            if shapes and isinstance(shapes, list):
                for shp in shapes:
                    try:
                        shp_type = shp.get('type', 'rect')
                        shp_text = shp.get('text', '')
                        shp_x = float(shp.get('x', 1))
                        shp_y = float(shp.get('y', 1))
                        shp_w = float(shp.get('w', 2))
                        shp_h = float(shp.get('h', 1))
                        fill_color = shp.get('fill_color', '')
                        text_color = shp.get('text_color', '')
                        font_size = int(shp.get('font_size', 12))

                        # Map type to MSO_SHAPE
                        shape_map = {
                            'rect': MSO_SHAPE.RECTANGLE,
                            'round_rect': MSO_SHAPE.ROUNDED_RECTANGLE,
                            'oval': MSO_SHAPE.OVAL,
                            'diamond': MSO_SHAPE.DIAMOND,
                            'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
                            'arrow': MSO_SHAPE.RIGHT_ARROW,
                            'star': MSO_SHAPE.STAR_5_POINT,
                            'heart': MSO_SHAPE.HEART,
                            'line': MSO_SHAPE.RECTANGLE,  # Use thin rect
                            'chevron': MSO_SHAPE.CHEVRON,
                            'hexagon': MSO_SHAPE.HEXAGON,
                            'pentagon': MSO_SHAPE.REGULAR_PENTAGON,
                        }
                        mso_type = shape_map.get(shp_type, MSO_SHAPE.RECTANGLE)

                        shape_obj = slide.shapes.add_shape(
                            mso_type,
                            Inches(shp_x), Inches(shp_y),
                            Inches(shp_w), Inches(shp_h)
                        )

                        if fill_color:
                            shape_obj.fill.solid()
                            shape_obj.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip('#'))

                        if shp_text:
                            tf = shape_obj.text_frame
                            tf.word_wrap = True
                            para = tf.paragraphs[0]
                            para.alignment = PP_ALIGN.CENTER
                            run = para.add_run()
                            run.text = shp_text
                            run.font.size = Pt(font_size)
                            if text_color:
                                run.font.color.rgb = RGBColor.from_string(text_color.lstrip('#'))
                    except Exception:
                        continue

            # Add speaker notes
            if notes:
                try:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes
                except Exception:
                    pass

        prs.save(str(out_file))
        size_kb = out_file.stat().st_size / 1024
        return f"PPTX created: {path}\nSlides: {len(slides)}\nSize: {size_kb:.1f} KB"

    except Exception as e:
        return f"PPTX create error: {e}\n{traceback.format_exc()}"


def edit_pptx(path: str, output: str = '', operations_json: str = '') -> str:
    """
    Edit an existing PowerPoint file with various operations.
    REAL editing — modifies actual .pptx file.

    operations_json format:
    {
      "add_slides": [
        {"title": "New Slide", "content": "Body text", "bullets": [...]}
      ],
      "delete_slides": [3, 5],
      "replace_text": {"old_text": "find this", "new_text": "replace with this"},
      "add_shapes": [
        {"slide": 1, "type": "rect", "x": 1, "y": 2, "w": 3, "h": 1,
         "text": "Hello", "fill_color": "#FF0000"}
      ],
      "add_notes": {"slide": 1, "notes": "Speaker notes here"},
      "reorder_slides": [3, 1, 2],  // new order (1-indexed)
      "set_background": {"slide": "all", "color": "#1a1a2e"},
      "add_image": {
        "slide": 1,
        "image_path": "/path/to/image.png",
        "x": 1, "y": 2, "w": 4, "h": 3
      }
    }
    """
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"

    try:
        ops = json.loads(operations_json) if operations_json else {}
    except json.JSONDecodeError as e:
        return f"Error: Invalid JSON in operations: {e}"

    if not ops:
        return "Error: No operations specified"

    try:
        prs = Presentation(str(p))
        results = []

        # Add slides
        add_slides = ops.get('add_slides', [])
        if add_slides:
            for slide_data in add_slides:
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                if slide_data.get('title'):
                    try:
                        slide.shapes.title.text_frame.text = slide_data['title']
                    except Exception:
                        pass
                if slide_data.get('content'):
                    try:
                        if len(slide.placeholders) > 1:
                            slide.placeholders[1].text_frame.text = slide_data['content']
                    except Exception:
                        pass
                bullets = slide_data.get('bullets', [])
                if bullets:
                    try:
                        if len(slide.placeholders) > 1:
                            tf = slide.placeholders[1].text_frame
                            tf.clear()
                            for b in bullets:
                                p_para = tf.add_paragraph()
                                p_para.add_run().text = str(b)
                    except Exception:
                        pass
                if slide_data.get('notes'):
                    try:
                        slide.notes_slide.notes_text_frame.text = slide_data['notes']
                    except Exception:
                        pass
                results.append(f"Added slide: {slide_data.get('title', '(untitled)')}")

        # Delete slides
        delete_slides = ops.get('delete_slides', [])
        if delete_slides:
            # Sort in reverse order to delete from end first
            for slide_num in sorted(set(delete_slides), reverse=True):
                idx = slide_num - 1
                if 0 <= idx < len(prs.slides):
                    rId = prs.slides._sldIdLst[idx].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[idx]
                    results.append(f"Deleted slide {slide_num}")
                else:
                    results.append(f"Slide {slide_num} not found (total: {len(prs.slides)})")

        # Replace text (across all slides)
        replace_text = ops.get('replace_text', {})
        if replace_text:
            old_text = replace_text.get('old_text', '')
            new_text = replace_text.get('new_text', '')
            if old_text:
                count = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if old_text in run.text:
                                        run.text = run.text.replace(old_text, new_text)
                                        count += 1
                results.append(f"Replaced '{old_text}' -> '{new_text}' ({count} occurrences)")

        # Add shapes to existing slides
        add_shapes = ops.get('add_shapes', [])
        if add_shapes:
            for shp in add_shapes:
                slide_idx = shp.get('slide', 1) - 1
                if 0 <= slide_idx < len(prs.slides):
                    slide = prs.slides[slide_idx]
                    try:
                        shape_map = {
                            'rect': MSO_SHAPE.RECTANGLE,
                            'round_rect': MSO_SHAPE.ROUNDED_RECTANGLE,
                            'oval': MSO_SHAPE.OVAL,
                            'diamond': MSO_SHAPE.DIAMOND,
                            'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
                            'arrow': MSO_SHAPE.RIGHT_ARROW,
                            'star': MSO_SHAPE.STAR_5_POINT,
                            'hexagon': MSO_SHAPE.HEXAGON,
                        }
                        mso_type = shape_map.get(shp.get('type', 'rect'), MSO_SHAPE.RECTANGLE)
                        shape_obj = slide.shapes.add_shape(
                            mso_type,
                            Inches(float(shp.get('x', 1))),
                            Inches(float(shp.get('y', 1))),
                            Inches(float(shp.get('w', 2))),
                            Inches(float(shp.get('h', 1)))
                        )
                        fill_color = shp.get('fill_color', '')
                        if fill_color:
                            shape_obj.fill.solid()
                            shape_obj.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip('#'))
                        shp_text = shp.get('text', '')
                        if shp_text:
                            tf = shape_obj.text_frame
                            tf.word_wrap = True
                            run = tf.paragraphs[0].add_run()
                            run.text = shp_text
                            fs = shp.get('font_size', 12)
                            run.font.size = Pt(int(fs))
                            text_color = shp.get('text_color', '')
                            if text_color:
                                run.font.color.rgb = RGBColor.from_string(text_color.lstrip('#'))
                        results.append(f"Added shape to slide {slide_idx + 1}")
                    except Exception as e:
                        results.append(f"Shape error on slide {slide_idx + 1}: {e}")

        # Add speaker notes
        add_notes = ops.get('add_notes', {})
        if add_notes:
            slide_num = add_notes.get('slide', 1)
            notes_text = add_notes.get('notes', '')
            idx = slide_num - 1
            if 0 <= idx < len(prs.slides) and notes_text:
                try:
                    prs.slides[idx].notes_slide.notes_text_frame.text = notes_text
                    results.append(f"Added notes to slide {slide_num}")
                except Exception as e:
                    results.append(f"Notes error: {e}")

        # Set background
        set_bg = ops.get('set_background', {})
        if set_bg:
            bg_color = set_bg.get('color', '')
            slide_target = set_bg.get('slide', 'all')
            if bg_color:
                try:
                    rgb = RGBColor.from_string(bg_color.lstrip('#'))
                    if slide_target == 'all':
                        for slide in prs.slides:
                            bg = slide.background
                            fill = bg.fill
                            fill.solid()
                            fill.fore_color.rgb = rgb
                        results.append(f"Set background color for all slides: {bg_color}")
                    else:
                        idx = int(slide_target) - 1
                        if 0 <= idx < len(prs.slides):
                            bg = prs.slides[idx].background
                            fill = bg.fill
                            fill.solid()
                            fill.fore_color.rgb = rgb
                            results.append(f"Set background for slide {slide_target}: {bg_color}")
                except Exception as e:
                    results.append(f"Background error: {e}")

        # Add image
        add_image = ops.get('add_image', {})
        if add_image:
            slide_num = add_image.get('slide', 1)
            img_path = add_image.get('image_path', '')
            idx = slide_num - 1
            if 0 <= idx < len(prs.slides) and img_path:
                img_p = Path(img_path).expanduser()
                if img_p.exists():
                    try:
                        prs.slides[idx].shapes.add_picture(
                            str(img_p),
                            Inches(float(add_image.get('x', 1))),
                            Inches(float(add_image.get('y', 1))),
                            Inches(float(add_image.get('w', 4))),
                            Inches(float(add_image.get('h', 3)))
                        )
                        results.append(f"Added image to slide {slide_num}")
                    except Exception as e:
                        results.append(f"Image error: {e}")
                else:
                    results.append(f"Image not found: {img_path}")

        # Save
        out_path = Path(output).expanduser() if output else p
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        size_kb = out_path.stat().st_size / 1024

        return f"PPTX edited: {out_path}\nSize: {size_kb:.1f} KB\nOperations:\n" + '\n'.join(f"  - {r}" for r in results)

    except Exception as e:
        return f"PPTX edit error: {e}\n{traceback.format_exc()}"


def pptx_info(path: str) -> str:
    """Get detailed information about a PPTX file — metadata, slide summary, media."""
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"

    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"

    try:
        prs = Presentation(str(p))
        lines = []
        lines.append(f"PPTX Info: {p.name}")
        lines.append("=" * 50)
        lines.append(f"File size: {p.stat().st_size / 1024:.1f} KB")
        lines.append(f"Total slides: {len(prs.slides)}")

        width_in = Emu(prs.slide_width).inches
        height_in = Emu(prs.slide_height).inches
        lines.append(f"Slide size: {width_in:.2f}\" x {height_in:.2f}\"")

        # Count slide layouts
        lines.append(f"Slide layouts available: {len(prs.slide_layouts)}")

        # Per-slide summary
        total_shapes = 0
        total_text_chars = 0
        total_images = 0
        total_tables = 0
        slides_with_notes = 0

        for i, slide in enumerate(prs.slides, 1):
            shape_count = len(slide.shapes)
            total_shapes += shape_count
            slide_text = ''
            img_count = 0
            tbl_count = 0
            has_title = False
            layout_name = ''

            if slide.slide_layout:
                layout_name = slide.slide_layout.name

            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text
                    slide_text += t
                    total_text_chars += len(t)
                if shape.has_table:
                    tbl_count += 1
                if shape.shape_type == 13:  # PICTURE
                    img_count += 1
                if shape.name == 'Title 1' or (hasattr(shape, 'is_placeholder') and shape.has_text_frame and shape == slide.shapes.title):
                    has_title = True

            total_images += img_count
            total_tables += tbl_count
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slides_with_notes += 1

            title_preview = slide_text[:80].replace('\n', ' ') if slide_text else '(empty)'
            lines.append(f"  Slide {i}: {shape_count} shapes | {layout_name} | \"{title_preview}\"")

        lines.append("")
        lines.append("Summary:")
        lines.append(f"  Total shapes: {total_shapes}")
        lines.append(f"  Total text characters: {total_text_chars}")
        lines.append(f"  Total images: {total_images}")
        lines.append(f"  Total tables: {total_tables}")
        lines.append(f"  Slides with notes: {slides_with_notes}")

        # Read core.xml for metadata
        try:
            with zipfile.ZipFile(str(p), 'r') as zf:
                if 'docProps/core.xml' in zf.namelist():
                    import xml.etree.ElementTree as ET
                    core_xml = zf.read('docProps/core.xml').decode('utf-8')
                    root = ET.fromstring(core_xml)
                    ns = {'cp': 'http://schemas.openxmlformats.org/package/2006/metadata/core-properties',
                          'dc': 'http://purl.org/dc/elements/1.1/',
                          'dcterms': 'http://purl.org/dc/terms/',
                          'xsi': 'http://www.w3.org/2001/XMLSchema-instance'}
                    lines.append("")
                    lines.append("Metadata:")
                    for tag, label in [
                        ('dc:title', 'Title'), ('dc:creator', 'Creator'),
                        ('dc:subject', 'Subject'), ('dc:description', 'Description'),
                        ('dc:language', 'Language'), ('cp:keywords', 'Keywords'),
                        ('cp:lastModifiedBy', 'Last Modified By'),
                    ]:
                        el = root.find(tag, ns)
                        if el is not None and el.text:
                            lines.append(f"  {label}: {el.text}")
                    for tag, label in [
                        ('dcterms:created', 'Created'), ('dcterms:modified', 'Modified'),
                    ]:
                        el = root.find(tag, ns)
                        if el is not None and el.text:
                            lines.append(f"  {label}: {el.text}")
        except Exception:
            pass

        return '\n'.join(lines)

    except Exception as e:
        return f"PPTX info error: {e}"


# ═══════════════════════════════════════════════════════════════
# XLSX TOOLS — Excel Spreadsheets (REAL)
# ═══════════════════════════════════════════════════════════════

def read_xlsx(path: str, sheet: str = '', max_rows: int = 100, header: bool = True) -> str:
    """
    Read an Excel (.xlsx) file and return content as formatted text.
    REAL reading with openpyxl.
    """
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"

    try:
        import openpyxl
    except ImportError:
        return "Error: openpyxl not installed. Run: pip install openpyxl"

    try:
        wb = openpyxl.load_workbook(str(p), data_only=True)
        lines = []
        lines.append(f"XLSX: {p.name}")
        lines.append(f"File size: {p.stat().st_size / 1024:.1f} KB")
        lines.append(f"Sheets: {wb.sheetnames}")

        # Core properties
        try:
            cp = wb.properties
            if cp.title or cp.creator or cp.subject:
                lines.append(f"Properties: title={cp.title or '-'} | creator={cp.creator or '-'} | subject={cp.subject or '-'}")
        except Exception:
            pass

        lines.append("=" * 60)

        # Read specified sheet or all
        sheets_to_read = [sheet] if sheet else wb.sheetnames
        for sname in sheets_to_read:
            if sname not in wb.sheetnames:
                lines.append(f"\nSheet '{sname}' not found. Available: {wb.sheetnames}")
                continue

            ws = wb[sname]
            lines.append(f"\n--- Sheet: {sname} ---")
            lines.append(f"Dimensions: {ws.dimensions}")
            lines.append(f"Rows: {ws.max_row} | Cols: {ws.max_column}")

            # Read data
            row_count = 0
            for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, max_rows), values_only=False):
                cells = []
                is_header = (row_count == 0 and header)
                for cell in row:
                    val = cell.value
                    if val is None:
                        cells.append('')
                    elif isinstance(val, (int, float)):
                        cells.append(str(val))
                    elif isinstance(val, datetime):
                        cells.append(val.strftime('%Y-%m-%d %H:%M:%S'))
                    else:
                        cells.append(str(val))

                row_text = ' | '.join(cells)
                if is_header:
                    lines.append(f"  {'─' * len(row_text)}")
                    lines.append(f"  {row_text}")
                    lines.append(f"  {'─' * len(row_text)}")
                else:
                    lines.append(f"  {row_text}")
                row_count += 1

            if ws.max_row > max_rows:
                lines.append(f"  ... ({ws.max_row - max_rows} more rows)")

            # Merged cells info
            if ws.merged_cells.ranges:
                lines.append(f"Merged cells: {len(ws.merged_cells.ranges)} ranges")
                for mc in list(ws.merged_cells.ranges)[:10]:
                    lines.append(f"  {mc}")

        wb.close()
        return '\n'.join(lines)

    except Exception as e:
        return f"XLSX read error: {e}\n{traceback.format_exc()}"


def create_xlsx(path: str, sheets_json: str = '', title: str = '') -> str:
    """
    Create an Excel (.xlsx) file with data.
    REAL creation with openpyxl — supports formulas, formatting, multiple sheets, headers.

    sheets_json format:
    {
      "sheets": [
        {
          "name": "Sheet1",
          "headers": ["Name", "Age", "Score"],
          "data": [
            ["Alice", 25, 95.5],
            ["Bob", 30, 87.0]
          ],
          "formulas": {
            "C4": "=SUM(C2:C3)",
            "D2": "=B2*2"
          },
          "column_widths": {"A": 20, "B": 10, "C": 10},
          "freeze_panes": "A2",
          "auto_filter": true,
          "bold_headers": true,
          "number_format": {"C": "0.00"},
          "header_color": "#4472C4",
          "header_text_color": "#FFFFFF"
        }
      ]
    }

    Simple mode (if sheets_json is empty but title is set):
    Creates a single-sheet workbook with the title.
    """
    p = Path(path).expanduser()
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        return "Error: openpyxl not installed. Run: pip install openpyxl"

    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        wb = openpyxl.Workbook()
        ws_default = wb.active
        ws_default.title = "Sheet1"

        if title and not sheets_json:
            # Simple mode
            ws_default['A1'] = title
            ws_default['A1'].font = Font(size=16, bold=True)
            wb.save(str(p))
            size_kb = p.stat().st_size / 1024
            return f"XLSX created: {path} (1 sheet, title only)\nSize: {size_kb:.1f} KB"

        sheets_data = json.loads(sheets_json) if sheets_json else {}

        if isinstance(sheets_data, list):
            sheets_list = sheets_data
        elif isinstance(sheets_data, dict):
            sheets_list = sheets_data.get('sheets', [])
        else:
            sheets_list = []

        if not sheets_list:
            wb.save(str(p))
            return f"XLSX created (empty): {path}"

        # Remove default sheet
        if len(sheets_list) > 0:
            wb.remove(ws_default)

        total_rows = 0
        for sheet_idx, sheet_data in enumerate(sheets_list):
            sheet_name = sheet_data.get('name', f'Sheet{sheet_idx + 1}')
            ws = wb.create_sheet(title=sheet_name[:31])  # Max 31 chars

            headers = sheet_data.get('headers', [])
            data = sheet_data.get('data', [])
            formulas = sheet_data.get('formulas', {})
            col_widths = sheet_data.get('column_widths', {})
            freeze_panes = sheet_data.get('freeze_panes', '')
            auto_filter = sheet_data.get('auto_filter', False)
            bold_headers = sheet_data.get('bold_headers', True)
            number_format = sheet_data.get('number_format', {})
            header_color = sheet_data.get('header_color', '')
            header_text_color = sheet_data.get('header_text_color', '')

            # Write headers
            if headers:
                for j, h in enumerate(headers):
                    cell = ws.cell(row=1, column=j + 1, value=h)
                    if bold_headers:
                        cell.font = Font(bold=True)
                    if header_color:
                        try:
                            cell.fill = PatternFill(
                                start_color=header_color.lstrip('#'),
                                end_color=header_color.lstrip('#'),
                                fill_type='solid'
                            )
                        except Exception:
                            pass
                    if header_text_color:
                        try:
                            cell.font = Font(bold=True, color=header_text_color.lstrip('#'))
                        except Exception:
                            pass

                # Add border to headers
                thin_border = Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                )
                for j in range(len(headers)):
                    ws.cell(row=1, column=j + 1).border = thin_border
                    ws.cell(row=1, column=j + 1).alignment = Alignment(horizontal='center')

            # Write data
            for i, row in enumerate(data):
                for j, val in enumerate(row):
                    cell = ws.cell(row=i + 2, column=j + 1)
                    if isinstance(val, str):
                        # Check if it's a number
                        try:
                            cell.value = int(val)
                        except ValueError:
                            try:
                                cell.value = float(val)
                            except ValueError:
                                cell.value = val
                    else:
                        cell.value = val

                total_rows += 1

            # Apply formulas
            if formulas:
                for cell_ref, formula in formulas.items():
                    ws[cell_ref] = formula

            # Apply number formats
            if number_format:
                for col_letter, fmt in number_format.items():
                    if col_letter.isalpha():
                        for row in range(2, len(data) + 2):
                            ws[f"{col_letter}{row}"].number_format = fmt

            # Set column widths
            if col_widths:
                for col_letter, width in col_widths.items():
                    ws.column_dimensions[col_letter].width = int(width)

            # Freeze panes
            if freeze_panes:
                ws.freeze_panes = freeze_panes

            # Auto filter
            if auto_filter and headers:
                max_col_letter = get_column_letter(len(headers))
                ws.auto_filter.ref = f"A1:{max_col_letter}{len(data) + 1}"

            total_rows += len(data)

        wb.save(str(p))
        size_kb = p.stat().st_size / 1024
        return f"XLSX created: {path}\nSheets: {len(sheets_list)}\nTotal rows: {total_rows}\nSize: {size_kb:.1f} KB"

    except Exception as e:
        return f"XLSX create error: {e}\n{traceback.format_exc()}"


def edit_xlsx(path: str, output: str = '', operations_json: str = '') -> str:
    """
    Edit an existing Excel file with various operations.
    REAL editing with openpyxl.

    operations_json format:
    {
      "add_sheet": {"name": "NewSheet", "data": [["a","b"],[1,2]]},
      "delete_sheet": "Sheet2",
      "rename_sheet": {"old": "Sheet1", "new": "Renamed"},
      "add_rows": {"sheet": "Sheet1", "data": [["new1","new2"]], "start_row": 10},
      "update_cell": {"sheet": "Sheet1", "cell": "A1", "value": "Updated"},
      "add_formula": {"sheet": "Sheet1", "cell": "C1", "formula": "=SUM(A1:B1)"},
      "set_column_width": {"sheet": "Sheet1", "column": "A", "width": 25},
      "set_row_height": {"sheet": "Sheet1", "row": 1, "height": 30},
      "delete_rows": {"sheet": "Sheet1", "rows": [5, 6, 7]},
      "delete_columns": {"sheet": "Sheet1", "columns": ["C", "D"]},
      "find_replace": {"sheet": "Sheet1", "find": "old", "replace": "new"},
      "sort_range": {"sheet": "Sheet1", "range": "A1:C10", "by_column": "A", "ascending": true},
      "merge_cells": {"sheet": "Sheet1", "range": "A1:B1"},
      "unmerge_cells": {"sheet": "Sheet1", "range": "A1:B1"},
      "add_chart": {
        "sheet": "Sheet1",
        "chart_type": "bar" | "line" | "pie",
        "title": "Sales Chart",
        "data_range": "A1:B10",
        "categories": "A2:A10",
        "values": "B2:B10",
        "position": "D2"
      },
      "set_cell_format": {
        "sheet": "Sheet1",
        "range": "A1:A10",
        "bold": true,
        "italic": false,
        "font_size": 14,
        "font_color": "#FF0000",
        "bg_color": "#FFFF00",
        "number_format": "0.00",
        "alignment": "center"
      }
    }
    """
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"

    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        from openpyxl.utils import get_column_letter
    except ImportError:
        return "Error: openpyxl not installed. Run: pip install openpyxl"

    try:
        ops = json.loads(operations_json) if operations_json else {}
    except json.JSONDecodeError as e:
        return f"Error: Invalid JSON: {e}"

    if not ops:
        return "Error: No operations specified"

    try:
        wb = openpyxl.load_workbook(str(p))
        results = []

        def _get_ws(name):
            if name in wb.sheetnames:
                return wb[name]
            results.append(f"Sheet '{name}' not found")
            return None

        # Add sheet
        add_sheet = ops.get('add_sheet', {})
        if add_sheet and isinstance(add_sheet, dict):
            sname = add_sheet.get('name', 'NewSheet')
            ws = wb.create_sheet(title=sname[:31])
            data = add_sheet.get('data', [])
            for i, row in enumerate(data):
                for j, val in enumerate(row):
                    ws.cell(row=i + 1, column=j + 1, value=val)
            results.append(f"Added sheet '{sname}' with {len(data)} rows")

        # Delete sheet
        del_sheet = ops.get('delete_sheet', '')
        if del_sheet:
            if del_sheet in wb.sheetnames:
                if len(wb.sheetnames) > 1:
                    wb.remove(wb[del_sheet])
                    results.append(f"Deleted sheet '{del_sheet}'")
                else:
                    results.append("Cannot delete the last sheet")
            else:
                results.append(f"Sheet '{del_sheet}' not found")

        # Rename sheet
        rename = ops.get('rename_sheet', {})
        if rename and isinstance(rename, dict):
            old_name = rename.get('old', '')
            new_name = rename.get('new', '')
            if old_name in wb.sheetnames and new_name:
                wb[old_name].title = new_name[:31]
                results.append(f"Renamed '{old_name}' -> '{new_name}'")

        # Add rows
        add_rows = ops.get('add_rows', {})
        if add_rows and isinstance(add_rows, dict):
            ws = _get_ws(add_rows.get('sheet', ''))
            if ws:
                data = add_rows.get('data', [])
                start_row = add_rows.get('start_row', ws.max_row + 1)
                for i, row in enumerate(data):
                    for j, val in enumerate(row):
                        ws.cell(row=start_row + i, column=j + 1, value=val)
                results.append(f"Added {len(data)} rows to '{add_rows.get('sheet')}' starting at row {start_row}")

        # Update cell
        update_cell = ops.get('update_cell', {})
        if update_cell and isinstance(update_cell, dict):
            ws = _get_ws(update_cell.get('sheet', ''))
            if ws:
                cell = update_cell.get('cell', 'A1')
                value = update_cell.get('value', '')
                ws[cell] = value
                results.append(f"Updated {cell} = '{value}'")

        # Add formula
        add_formula = ops.get('add_formula', {})
        if add_formula and isinstance(add_formula, dict):
            ws = _get_ws(add_formula.get('sheet', ''))
            if ws:
                cell = add_formula.get('cell', 'A1')
                formula = add_formula.get('formula', '')
                ws[cell] = formula
                results.append(f"Set formula {cell} = {formula}")

        # Column width
        set_col_width = ops.get('set_column_width', {})
        if set_col_width and isinstance(set_col_width, dict):
            ws = _get_ws(set_col_width.get('sheet', ''))
            if ws:
                col = set_col_width.get('column', 'A')
                width = set_col_width.get('width', 15)
                ws.column_dimensions[col].width = int(width)
                results.append(f"Column {col} width = {width}")

        # Row height
        set_row_height = ops.get('set_row_height', {})
        if set_row_height and isinstance(set_row_height, dict):
            ws = _get_ws(set_row_height.get('sheet', ''))
            if ws:
                row = int(set_row_height.get('row', 1))
                height = set_row_height.get('height', 15)
                ws.row_dimensions[row].height = float(height)
                results.append(f"Row {row} height = {height}")

        # Delete rows
        del_rows = ops.get('delete_rows', {})
        if del_rows and isinstance(del_rows, dict):
            ws = _get_ws(del_rows.get('sheet', ''))
            if ws:
                rows = del_rows.get('rows', [])
                for r in sorted(rows, reverse=True):
                    ws.delete_rows(r)
                results.append(f"Deleted {len(rows)} rows from '{del_rows.get('sheet')}'")

        # Delete columns
        del_cols = ops.get('delete_columns', {})
        if del_cols and isinstance(del_cols, dict):
            ws = _get_ws(del_cols.get('sheet', ''))
            if ws:
                cols = del_cols.get('columns', [])
                for c in sorted(cols, reverse=True):
                    ws.delete_cols(c)
                results.append(f"Deleted columns {cols} from '{del_cols.get('sheet')}'")

        # Find and replace
        find_replace = ops.get('find_replace', {})
        if find_replace and isinstance(find_replace, dict):
            ws = _get_ws(find_replace.get('sheet', ''))
            if ws:
                find_text = find_replace.get('find', '')
                replace_text = find_replace.get('replace', '')
                count = 0
                for row in ws.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and find_text in cell.value:
                            cell.value = cell.value.replace(find_text, replace_text)
                            count += 1
                results.append(f"Replaced '{find_text}' -> '{replace_text}' ({count} cells)")

        # Sort range
        sort_range = ops.get('sort_range', {})
        if sort_range and isinstance(sort_range, dict):
            ws = _get_ws(sort_range.get('sheet', ''))
            if ws:
                range_str = sort_range.get('range', 'A1:C10')
                by_col = sort_range.get('by_column', 'A')
                ascending = sort_range.get('ascending', True)
                try:
                    # Parse range
                    match = re.match(r'([A-Z]+)(\d+):([A-Z]+)(\d+)', range_str)
                    if match:
                        start_col, start_row, end_col, end_row = match.groups()
                        col_idx = ord(by_col.upper()) - ord('A')
                        rows_data = []
                        for row in ws.iter_rows(min_row=int(start_row), max_row=int(end_row),
                                               min_col=ord(start_col) - ord('A') + 1,
                                               max_col=ord(end_col) - ord('A') + 1):
                            rows_data.append([cell.value for cell in row])
                        rows_data.sort(key=lambda r: (r[col_idx] is None, r[col_idx]),
                                       reverse=not ascending)
                        for i, rd in enumerate(rows_data):
                            for j, val in enumerate(rd):
                                ws.cell(row=int(start_row) + i,
                                        column=ord(start_col) - ord('A') + 1 + j).value = val
                        results.append(f"Sorted {range_str} by column {by_col}")
                except Exception as e:
                    results.append(f"Sort error: {e}")

        # Merge cells
        merge_cells = ops.get('merge_cells', {})
        if merge_cells and isinstance(merge_cells, dict):
            ws = _get_ws(merge_cells.get('sheet', ''))
            if ws:
                range_str = merge_cells.get('range', 'A1:B1')
                try:
                    ws.merge_cells(range_str)
                    results.append(f"Merged cells {range_str}")
                except Exception as e:
                    results.append(f"Merge error: {e}")

        # Unmerge cells
        unmerge_cells = ops.get('unmerge_cells', {})
        if unmerge_cells and isinstance(unmerge_cells, dict):
            ws = _get_ws(unmerge_cells.get('sheet', ''))
            if ws:
                range_str = unmerge_cells.get('range', 'A1:B1')
                try:
                    ws.unmerge_cells(range_str)
                    results.append(f"Unmerged cells {range_str}")
                except Exception as e:
                    results.append(f"Unmerge error: {e}")

        # Add chart
        add_chart = ops.get('add_chart', {})
        if add_chart and isinstance(add_chart, dict):
            ws = _get_ws(add_chart.get('sheet', ''))
            if ws:
                try:
                    chart_type = add_chart.get('chart_type', 'bar')
                    chart_title = add_chart.get('title', 'Chart')
                    data_range = add_chart.get('data_range', '')
                    categories = add_chart.get('categories', '')
                    values = add_chart.get('values', '')
                    position = add_chart.get('position', 'E2')

                    if chart_type == 'bar':
                        chart = BarChart()
                    elif chart_type == 'line':
                        chart = LineChart()
                    elif chart_type == 'pie':
                        chart = PieChart()
                    else:
                        chart = BarChart()

                    chart.title = chart_title
                    chart.style = 10

                    if data_range:
                        # Reference requires sheetname!range format or we can use min/max
                        dr = data_range.split('!')[-1].strip()
                        try:
                            ref = Reference(ws, range_string=f"{ws.title}!{dr}")
                        except Exception:
                            # Fallback: parse range manually
                            try:
                                match = re.match(r'([A-Z]+)(\d+):([A-Z]+)(\d+)', dr)
                                if match:
                                    from openpyxl.utils import column_index_from_string
                                    ref = Reference(ws,
                                        min_col=column_index_from_string(match.group(1)),
                                        min_row=int(match.group(2)),
                                        max_col=column_index_from_string(match.group(3)),
                                        max_row=int(match.group(4)))
                                else:
                                    ref = Reference(ws, range_string=dr)
                            except Exception:
                                ref = Reference(ws, range_string=dr)
                        chart.add_data(ref, titles_from_data=True)
                    if categories:
                        cr = categories.split('!')[-1].strip()
                        try:
                            cats = Reference(ws, range_string=f"{ws.title}!{cr}")
                        except Exception:
                            cats = Reference(ws, range_string=cr)
                        chart.set_categories(cats)

                    chart.width = 15
                    chart.height = 10
                    ws.add_chart(chart, position)
                    results.append(f"Added {chart_type} chart at {position}")
                except Exception as e:
                    results.append(f"Chart error: {e}")

        # Set cell formatting
        set_format = ops.get('set_cell_format', {})
        if set_format and isinstance(set_format, dict):
            ws = _get_ws(set_format.get('sheet', ''))
            if ws:
                range_str = set_format.get('range', 'A1:A10')
                font_kwargs = {}
                if set_format.get('bold'):
                    font_kwargs['bold'] = True
                if set_format.get('italic'):
                    font_kwargs['italic'] = True
                if set_format.get('font_size'):
                    font_kwargs['size'] = int(set_format['font_size'])
                if set_format.get('font_color'):
                    font_kwargs['color'] = set_format['font_color'].lstrip('#')

                bg_color = set_format.get('bg_color', '')
                num_fmt = set_format.get('number_format', '')
                alignment = set_format.get('alignment', '')

                try:
                    for row in ws[range_str]:
                        if isinstance(row, tuple):
                            cells = row
                        else:
                            cells = [row]
                        for cell in cells:
                            if font_kwargs:
                                cell.font = Font(**font_kwargs)
                            if bg_color:
                                cell.fill = PatternFill(
                                    start_color=bg_color.lstrip('#'),
                                    end_color=bg_color.lstrip('#'),
                                    fill_type='solid'
                                )
                            if num_fmt:
                                cell.number_format = num_fmt
                            if alignment:
                                align_map = {
                                    'center': 'center',
                                    'left': 'left',
                                    'right': 'right',
                                }
                                cell.alignment = Alignment(horizontal=align_map.get(alignment, 'left'))
                    results.append(f"Formatted range {range_str}")
                except Exception as e:
                    results.append(f"Format error: {e}")

        # Save
        out_path = Path(output).expanduser() if output else p
        out_path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(out_path))
        size_kb = out_path.stat().st_size / 1024

        return f"XLSX edited: {out_path}\nSize: {size_kb:.1f} KB\nOperations:\n" + '\n'.join(f"  - {r}" for r in results)

    except Exception as e:
        return f"XLSX edit error: {e}\n{traceback.format_exc()}"


def xlsx_info(path: str) -> str:
    """Get detailed info about an XLSX file — sheets, data stats, formulas, etc."""
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"

    try:
        import openpyxl
        from openpyxl.utils import get_column_letter
    except ImportError:
        return "Error: openpyxl not installed. Run: pip install openpyxl"

    try:
        wb = openpyxl.load_workbook(str(p))
        lines = []
        lines.append(f"XLSX Info: {p.name}")
        lines.append("=" * 50)
        lines.append(f"File size: {p.stat().st_size / 1024:.1f} KB")
        lines.append(f"Sheets: {len(wb.sheetnames)}")

        # Properties
        try:
            cp = wb.properties
            lines.append("")
            lines.append("Properties:")
            lines.append(f"  Title: {cp.title or '(none)'}")
            lines.append(f"  Creator: {cp.creator or '(none)'}")
            lines.append(f"  Subject: {cp.subject or '(none)'}")
            lines.append(f"  Keywords: {cp.keywords or '(none)'}")
            lines.append(f"  Created: {cp.created}")
            lines.append(f"  Modified: {cp.modified}")
            lines.append(f"  Last Modified By: {cp.lastModifiedBy or '(none)'}")
        except Exception:
            pass

        # Per-sheet info
        for sname in wb.sheetnames:
            ws = wb[sname]
            lines.append(f"\nSheet: {sname}")
            lines.append(f"  Dimensions: {ws.dimensions}")
            lines.append(f"  Max row: {ws.max_row} | Max col: {ws.max_column}")

            # Count data types
            num_count = 0
            str_count = 0
            formula_count = 0
            empty_count = 0
            bool_count = 0
            date_count = 0

            for row in ws.iter_rows(max_row=min(ws.max_row, 5000)):
                for cell in row:
                    if cell.value is None:
                        empty_count += 1
                    elif isinstance(cell.value, bool):
                        bool_count += 1
                    elif isinstance(cell.value, (int, float)):
                        num_count += 1
                    elif isinstance(cell.value, datetime):
                        date_count += 1
                    elif isinstance(cell.value, str) and cell.value.startswith('='):
                        formula_count += 1
                    else:
                        str_count += 1

            lines.append(f"  Numbers: {num_count} | Strings: {str_count} | Dates: {date_count}")
            lines.append(f"  Formulas: {formula_count} | Booleans: {bool_count} | Empty: {empty_count}")

            # Merged cells
            mc = len(ws.merged_cells.ranges)
            if mc > 0:
                lines.append(f"  Merged cells: {mc}")

            # Freeze panes
            if ws.freeze_panes:
                lines.append(f"  Freeze panes: {ws.freeze_panes}")

            # Auto filter
            if ws.auto_filter and ws.auto_filter.ref:
                lines.append(f"  Auto filter: {ws.auto_filter.ref}")

            # Column widths
            widths = {col: ws.column_dimensions[col].width for col in ws.column_dimensions if ws.column_dimensions[col].width}
            if widths:
                lines.append(f"  Column widths: {widths}")

        wb.close()
        return '\n'.join(lines)

    except Exception as e:
        return f"XLSX info error: {e}"


# ═══════════════════════════════════════════════════════════════
# DOCX EDIT TOOLS — Enhanced Word Editing (REAL)
# ═══════════════════════════════════════════════════════════════

def edit_docx(path: str, output: str = '', operations_json: str = '') -> str:
    """
    Edit an existing Word (.docx) file with various operations.
    REAL editing with python-docx.

    operations_json format:
    {
      "replace_text": {"find": "old text", "replace": "new text"},
      "replace_all": true,
      "append_text": "Text to append at the end",
      "append_content": "# New Section\\nNew paragraph here",
      "insert_heading": {"text": "New Heading", "level": 1, "position": "end"},
      "add_table": {
        "headers": ["Col1", "Col2", "Col3"],
        "rows": [["a", "b", "c"], ["d", "e", "f"]],
        "position": "end"
      },
      "set_properties": {
        "title": "Document Title",
        "author": "Author Name",
        "subject": "Subject",
        "keywords": "key1, key2",
        "category": "Reports"
      },
      "delete_paragraphs": {"index_start": 5, "index_end": 10},
      "add_page_break": {"position": "end"},
      "set_font": {
        "paragraph_index": 0,
        "bold": true,
        "italic": false,
        "size": 14,
        "color": "FF0000",
        "font_name": "Arial"
      }
    }
    """
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"

    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return "Error: python-docx not installed. Run: pip install python-docx"

    try:
        ops = json.loads(operations_json) if operations_json else {}
    except json.JSONDecodeError as e:
        return f"Error: Invalid JSON: {e}"

    if not ops:
        return "Error: No operations specified"

    try:
        doc = Document(str(p))
        results = []

        # Replace text
        replace = ops.get('replace_text', {})
        if replace and isinstance(replace, dict):
            find_text = replace.get('find', '')
            replace_text = replace.get('replace', '')
            replace_all = ops.get('replace_all', True)
            if find_text:
                count = 0
                for para in doc.paragraphs:
                    if find_text in para.text:
                        if replace_all:
                            for run in para.runs:
                                if find_text in run.text:
                                    run.text = run.text.replace(find_text, replace_text)
                                    count += 1
                        else:
                            for run in para.runs:
                                if find_text in run.text:
                                    run.text = run.text.replace(find_text, replace_text)
                                    count += 1
                                    break
                            break
                results.append(f"Replaced '{find_text[:30]}' -> '{replace_text[:30]}' ({count} times)")

        # Also replace in tables
        if replace and isinstance(replace, dict) and replace.get('find', ''):
            find_text = replace.get('find', '')
            replace_text = replace.get('replace', '')
            if find_text:
                t_count = 0
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    if find_text in run.text:
                                        run.text = run.text.replace(find_text, replace_text)
                                        t_count += 1
                if t_count > 0:
                    results.append(f"Also replaced in tables: {t_count} occurrences")

        # Append text
        append_text = ops.get('append_text', '')
        if append_text:
            doc.add_paragraph(append_text)
            results.append(f"Appended text: '{append_text[:50]}...'")

        # Append content (markdown-like)
        append_content = ops.get('append_content', '')
        if append_content:
            lines = append_content.split('\n')
            for line in lines:
                stripped = line.strip()
                if not stripped:
                    doc.add_paragraph('')
                elif stripped.startswith('#### '):
                    doc.add_heading(stripped[5:], level=4)
                elif stripped.startswith('### '):
                    doc.add_heading(stripped[4:], level=3)
                elif stripped.startswith('## '):
                    doc.add_heading(stripped[3:], level=2)
                elif stripped.startswith('# '):
                    doc.add_heading(stripped[2:], level=1)
                elif stripped.startswith('- ') or stripped.startswith('* '):
                    doc.add_paragraph(stripped[2:], style='List Bullet')
                elif re.match(r'^\d+\.\s', stripped):
                    doc.add_paragraph(stripped, style='List Number')
                else:
                    doc.add_paragraph(stripped)
            results.append(f"Appended content: {len(lines)} lines")

        # Insert heading
        insert_heading = ops.get('insert_heading', {})
        if insert_heading and isinstance(insert_heading, dict):
            heading_text = insert_heading.get('text', '')
            level = insert_heading.get('level', 1)
            if heading_text:
                doc.add_heading(heading_text, level=level)
                results.append(f"Added heading: '{heading_text}' (level {level})")

        # Add table
        add_table = ops.get('add_table', {})
        if add_table and isinstance(add_table, dict):
            headers = add_table.get('headers', [])
            rows = add_table.get('rows', [])
            if headers:
                table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                table.style = 'Table Grid'
                for j, h in enumerate(headers):
                    table.rows[0].cells[j].text = h
                    # Bold headers
                    for paragraph in table.rows[0].cells[j].paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                for i, row in enumerate(rows):
                    for j, val in enumerate(row):
                        if j < len(headers):
                            table.rows[i + 1].cells[j].text = str(val)
                results.append(f"Added table: {len(headers)} cols x {len(rows)} rows")

        # Set properties
        set_props = ops.get('set_properties', {})
        if set_props and isinstance(set_props, dict):
            cp = doc.core_properties
            for prop, label in [
                ('title', 'Title'), ('author', 'Author'),
                ('subject', 'Subject'), ('keywords', 'Keywords'),
                ('category', 'Category'),
            ]:
                if prop in set_props:
                    setattr(cp, prop, set_props[prop])
                    results.append(f"Set property {label}: {set_props[prop]}")

        # Delete paragraphs
        del_paras = ops.get('delete_paragraphs', {})
        if del_paras and isinstance(del_paras, dict):
            idx_start = del_paras.get('index_start', 0)
            idx_end = del_paras.get('index_end', idx_start + 1)
            # We need to delete XML elements
            try:
                body = doc.element.body
                paras = list(body)
                to_remove = []
                para_count = 0
                for child in body:
                    if child.tag.endswith('}p'):
                        if idx_start <= para_count < idx_end:
                            to_remove.append(child)
                        para_count += 1
                for el in to_remove:
                    body.remove(el)
                results.append(f"Deleted paragraphs {idx_start} to {idx_end - 1}")
            except Exception as e:
                results.append(f"Delete paragraphs error: {e}")

        # Add page break
        add_pb = ops.get('add_page_break', {})
        if add_pb:
            doc.add_page_break()
            results.append("Added page break")

        # Set font
        set_font = ops.get('set_font', {})
        if set_font and isinstance(set_font, dict):
            para_idx = set_font.get('paragraph_index', 0)
            if 0 <= para_idx < len(doc.paragraphs):
                para = doc.paragraphs[para_idx]
                for run in para.runs:
                    if set_font.get('bold') is not None:
                        run.font.bold = set_font['bold']
                    if set_font.get('italic') is not None:
                        run.font.italic = set_font['italic']
                    if set_font.get('size'):
                        run.font.size = Pt(int(set_font['size']))
                    if set_font.get('color'):
                        run.font.color.rgb = RGBColor.from_string(set_font['color'].lstrip('#'))
                    if set_font.get('font_name'):
                        run.font.name = set_font['font_name']
                results.append(f"Set font on paragraph {para_idx}")

        # Save
        out_path = Path(output).expanduser() if output else p
        out_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(out_path))
        size_kb = out_path.stat().st_size / 1024

        return f"DOCX edited: {out_path}\nSize: {size_kb:.1f} KB\nOperations:\n" + '\n'.join(f"  - {r}" for r in results)

    except Exception as e:
        return f"DOCX edit error: {e}\n{traceback.format_exc()}"


# ═══════════════════════════════════════════════════════════════
# CSV TOOLS (REAL)
# ═══════════════════════════════════════════════════════════════

def read_csv(path: str, delimiter: str = ',', max_rows: int = 200, encoding: str = '') -> str:
    """Read a CSV file and return formatted content."""
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"

    encodings = [encoding] if encoding else ['utf-8', 'latin-1', 'cp1252', 'utf-8-sig']

    for enc in encodings:
        try:
            with open(str(p), 'r', encoding=enc, newline='') as f:
                reader = csv.reader(f, delimiter=delimiter)
                rows = list(reader)

            lines = []
            lines.append(f"CSV: {p.name}")
            lines.append(f"File size: {p.stat().st_size / 1024:.1f} KB")
            lines.append(f"Rows: {len(rows)} | Delimiter: '{delimiter}' | Encoding: {enc}")

            if rows:
                lines.append(f"Columns: {len(rows[0])}")

            lines.append("=" * 60)

            display_rows = min(len(rows), max_rows)
            for i, row in enumerate(rows[:display_rows]):
                line = ' | '.join(str(cell) for cell in row)
                if i == 0:
                    lines.append(f"  {'─' * len(line)}")
                    lines.append(f"  {line}")
                    lines.append(f"  {'─' * len(line)}")
                else:
                    lines.append(f"  {line}")

            if len(rows) > max_rows:
                lines.append(f"  ... ({len(rows) - max_rows} more rows)")

            return '\n'.join(lines)

        except UnicodeDecodeError:
            continue
        except Exception as e:
            return f"CSV read error: {e}"

    return f"Error: Could not decode file with any encoding"


def create_csv(path: str, data_json: str = '', headers_json: str = '', delimiter: str = ',', encoding: str = 'utf-8') -> str:
    """
    Create a CSV file.
    data_json: JSON array of arrays: [["val1","val2"],["val3","val4"]]
    headers_json: JSON array: ["Col1","Col2"]
    """
    p = Path(path).expanduser()

    try:
        p.parent.mkdir(parents=True, exist_ok=True)

        headers = json.loads(headers_json) if headers_json else []
        data = json.loads(data_json) if data_json else []

        with open(str(p), 'w', encoding=encoding, newline='') as f:
            writer = csv.writer(f, delimiter=delimiter)
            if headers:
                writer.writerow(headers)
            for row in data:
                writer.writerow(row)

        size_kb = p.stat().st_size / 1024
        return f"CSV created: {path}\nRows: {len(data)} (+ {'1 header' if headers else 'no header'})\nSize: {size_kb:.1f} KB"

    except json.JSONDecodeError as e:
        return f"Error: Invalid JSON: {e}"
    except Exception as e:
        return f"CSV create error: {e}"


# ═══════════════════════════════════════════════════════════════
# DOCUMENT CONVERSION TOOLS (REAL)
# ═══════════════════════════════════════════════════════════════

def convert_document(input_path: str, output_path: str, options: str = '') -> str:
    """
    Convert between document formats.
    Supported conversions:
    - XLSX -> CSV
    - CSV -> XLSX
    - DOCX -> PDF (requires reportlab or pdfkit)
    - PPTX -> PDF (via LibreOffice if available)
    - Markdown -> DOCX
    - Markdown -> PDF
    - Text -> DOCX
    """
    inp = Path(input_path).expanduser()
    out = Path(output_path).expanduser()

    if not inp.exists():
        return f"Error: Input file not found: {input_path}"

    ext_in = inp.suffix.lower()
    ext_out = out.suffix.lower()

    out.parent.mkdir(parents=True, exist_ok=True)

    try:
        opts = json.loads(options) if options else {}
    except Exception:
        opts = {}

    # XLSX -> CSV
    if ext_in == '.xlsx' and ext_out == '.csv':
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(inp), data_only=True)
            ws = wb.active
            with open(str(out), 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    writer.writerow([str(v) if v is not None else '' for v in row])
            wb.close()
            return f"Converted: {input_path} -> {output_path}\nSize: {out.stat().st_size / 1024:.1f} KB"
        except ImportError:
            return "Error: openpyxl not installed. Run: pip install openpyxl"
        except Exception as e:
            return f"XLSX->CSV error: {e}"

    # CSV -> XLSX
    elif ext_in == '.csv' and ext_out == '.xlsx':
        try:
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = inp.stem[:31]

            encodings = ['utf-8', 'latin-1', 'utf-8-sig']
            for enc in encodings:
                try:
                    with open(str(inp), 'r', encoding=enc, newline='') as f:
                        reader = csv.reader(f)
                        for row_idx, row in enumerate(reader, 1):
                            for col_idx, val in enumerate(row, 1):
                                try:
                                    ws.cell(row=row_idx, column=col_idx, value=int(val))
                                except ValueError:
                                    try:
                                        ws.cell(row=row_idx, column=col_idx, value=float(val))
                                    except ValueError:
                                        ws.cell(row=row_idx, column=col_idx, value=val)
                    break
                except UnicodeDecodeError:
                    continue

            wb.save(str(out))
            return f"Converted: {input_path} -> {output_path}\nSize: {out.stat().st_size / 1024:.1f} KB"
        except ImportError:
            return "Error: openpyxl not installed"
        except Exception as e:
            return f"CSV->XLSX error: {e}"

    # Markdown -> DOCX
    elif ext_in == '.md' and ext_out == '.docx':
        try:
            from docx import Document
            doc = Document()
            content = inp.read_text(encoding='utf-8')
            lines = content.split('\n')
            for line in lines:
                s = line.strip()
                if not s:
                    doc.add_paragraph('')
                elif s.startswith('###### '):
                    doc.add_heading(s[7:], level=6)
                elif s.startswith('##### '):
                    doc.add_heading(s[6:], level=5)
                elif s.startswith('#### '):
                    doc.add_heading(s[5:], level=4)
                elif s.startswith('### '):
                    doc.add_heading(s[4:], level=3)
                elif s.startswith('## '):
                    doc.add_heading(s[3:], level=2)
                elif s.startswith('# '):
                    doc.add_heading(s[2:], level=1)
                elif s.startswith('- ') or s.startswith('* '):
                    doc.add_paragraph(s[2:], style='List Bullet')
                elif re.match(r'^\d+\.\s', s):
                    doc.add_paragraph(s, style='List Number')
                elif s.startswith('> '):
                    doc.add_paragraph(s[2:], style='Quote')
                elif s.startswith('---'):
                    doc.add_paragraph('_' * 40)
                elif s.startswith('```'):
                    continue
                else:
                    doc.add_paragraph(s)
            doc.save(str(out))
            return f"Converted: {input_path} -> {output_path}\nSize: {out.stat().st_size / 1024:.1f} KB"
        except ImportError:
            return "Error: python-docx not installed"
        except Exception as e:
            return f"MD->DOCX error: {e}"

    # Markdown -> PDF
    elif ext_in == '.md' and ext_out == '.pdf':
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        except ImportError:
            return "Error: reportlab not installed"
        try:
            doc = SimpleDocTemplate(str(out), pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            content = inp.read_text(encoding='utf-8')
            lines = content.split('\n')
            for line in lines:
                s = line.strip()
                if not s:
                    story.append(Spacer(1, 6))
                elif s.startswith('#### '):
                    story.append(Paragraph(s[5:], styles['Heading4']))
                elif s.startswith('### '):
                    story.append(Paragraph(s[4:], styles['Heading3']))
                elif s.startswith('## '):
                    story.append(Paragraph(s[3:], styles['Heading2']))
                elif s.startswith('# '):
                    story.append(Paragraph(s[2:], styles['Heading1']))
                elif s.startswith('- ') or s.startswith('* '):
                    story.append(Paragraph(f"  - {s[2:]}", styles['Normal']))
                else:
                    story.append(Paragraph(s, styles['Normal']))
            doc.build(story)
            return f"Converted: {input_path} -> {output_path}\nSize: {out.stat().st_size / 1024:.1f} KB"
        except Exception as e:
            return f"MD->PDF error: {e}"

    # Text -> DOCX
    elif ext_in == '.txt' and ext_out == '.docx':
        try:
            from docx import Document
            doc = Document()
            content = inp.read_text(encoding='utf-8')
            for para in content.split('\n'):
                doc.add_paragraph(para)
            doc.save(str(out))
            return f"Converted: {input_path} -> {output_path}\nSize: {out.stat().st_size / 1024:.1f} KB"
        except ImportError:
            return "Error: python-docx not installed"
        except Exception as e:
            return f"TXT->DOCX error: {e}"

    # JSON -> XLSX
    elif ext_in == '.json' and ext_out == '.xlsx':
        try:
            import openpyxl
            from openpyxl.styles import Font
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = inp.stem[:31]
            data = json.loads(inp.read_text(encoding='utf-8'))

            if isinstance(data, list) and data:
                if isinstance(data[0], dict):
                    # Array of objects -> table
                    headers = list(data[0].keys())
                    for j, h in enumerate(headers):
                        ws.cell(row=1, column=j + 1, value=h)
                        ws.cell(row=1, column=j + 1).font = Font(bold=True)
                    for i, item in enumerate(data):
                        for j, key in enumerate(headers):
                            ws.cell(row=i + 2, column=j + 1, value=item.get(key, ''))
                else:
                    # Simple array
                    for i, val in enumerate(data):
                        ws.cell(row=i + 1, column=1, value=val)
            elif isinstance(data, dict):
                # Key-value pairs
                ws.cell(row=1, column=1, value='Key').font = Font(bold=True)
                ws.cell(row=1, column=2, value='Value').font = Font(bold=True)
                for i, (k, v) in enumerate(data.items()):
                    ws.cell(row=i + 2, column=1, value=k)
                    ws.cell(row=i + 2, column=2, value=str(v))

            wb.save(str(out))
            return f"Converted: {input_path} -> {output_path}\nSize: {out.stat().st_size / 1024:.1f} KB"
        except ImportError:
            return "Error: openpyxl not installed"
        except Exception as e:
            return f"JSON->XLSX error: {e}"

    # JSON -> CSV
    elif ext_in == '.json' and ext_out == '.csv':
        try:
            data = json.loads(inp.read_text(encoding='utf-8'))
            with open(str(out), 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                if isinstance(data, list) and data and isinstance(data[0], dict):
                    headers = list(data[0].keys())
                    writer.writerow(headers)
                    for item in data:
                        writer.writerow([item.get(h, '') for h in headers])
                elif isinstance(data, dict):
                    writer.writerow(['Key', 'Value'])
                    for k, v in data.items():
                        writer.writerow([k, str(v)])
                else:
                    writer.writerow(['Value'])
                    for item in (data if isinstance(data, list) else [data]):
                        writer.writerow([str(item)])
            return f"Converted: {input_path} -> {output_path}\nSize: {out.stat().st_size / 1024:.1f} KB"
        except Exception as e:
            return f"JSON->CSV error: {e}"

    else:
        return (
            f"Error: Conversion from '{ext_in}' to '{ext_out}' is not directly supported.\n"
            f"Supported conversions:\n"
            f"  .xlsx -> .csv | .csv -> .xlsx\n"
            f"  .md -> .docx | .md -> .pdf\n"
            f"  .txt -> .docx\n"
            f"  .json -> .xlsx | .json -> .csv\n"
            f"\nFor other conversions, try LibreOffice command line:\n"
            f"  libreoffice --headless --convert-to pdf --outdir /output/ input.pptx"
        )


# ═══════════════════════════════════════════════════════════════
# TOOL DEFINITIONS FOR REGISTRATION
# Returns OpenAI function-calling format for all doc tools
# ═══════════════════════════════════════════════════════════════

def get_doc_tool_definitions() -> list[dict]:
    """Return all document tool definitions in OpenAI function-calling format."""
    return [
        # ── PPTX TOOLS ──
        {
            "type": "function",
            "function": {
                "name": "read_pptx",
                "description": (
                    "Read a PowerPoint (.pptx) file and extract all content including slide text, "
                    "tables, shapes, formatting details, speaker notes, and image info. "
                    "Returns detailed per-slide breakdown."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to the .pptx file"}
                    },
                    "required": ["path"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "create_pptx",
                "description": (
                    "Create a PowerPoint (.pptx) presentation from scratch. Supports multiple slides "
                    "with titles, content, bullets, tables, custom shapes, backgrounds, speaker notes. "
                    "Use slides_json to define all slides. Can create professional presentations with "
                    "colored shapes, formatted text, and data tables."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Output .pptx file path"},
                        "title": {"type": "string", "description": "Presentation title (simple mode)"},
                        "slides_json": {
                            "type": "string",
                            "description": (
                                "JSON array of slide objects. Each slide: {"
                                "\"layout\": \"title_slide\"|\"title_content\"|\"section_header\"|\"blank\", "
                                "\"title\": \"...\", \"subtitle\": \"...\", \"content\": \"body text\", "
                                "\"bullets\": [\"Point 1\", \"Point 2\"], "
                                "\"table\": {\"headers\": [...], \"rows\": [[...]]}, "
                                "\"notes\": \"speaker notes\", "
                                "\"background_color\": \"#HEX\", "
                                "\"title_color\": \"#HEX\", \"title_size\": 44, \"content_size\": 18, "
                                "\"shapes\": [{\"type\": \"rect\"|\"oval\"|\"star\"|\"arrow\"|\"diamond\"|\"hexagon\"|\"triangle\", "
                                "\"text\": \"label\", \"x\": 1, \"y\": 2, \"w\": 3, \"h\": 1, "
                                "\"fill_color\": \"#FF0000\", \"text_color\": \"#FFFFFF\", \"font_size\": 14}]}"
                            )
                        }
                    },
                    "required": ["path"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "edit_pptx",
                "description": (
                    "Edit an existing PowerPoint (.pptx) file. Can add/delete slides, replace text, "
                    "add shapes and images, set backgrounds, add speaker notes, and reorder slides. "
                    "Multiple operations in one call."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to existing .pptx file"},
                        "output": {"type": "string", "description": "Output path (default: overwrite input)"},
                        "operations_json": {
                            "type": "string",
                            "description": (
                                "JSON object with operations:\n"
                                "{\"add_slides\": [{\"title\": \"...\", \"content\": \"...\", \"bullets\": [...]}], "
                                "\"delete_slides\": [3, 5], "
                                "\"replace_text\": {\"old_text\": \"find\", \"new_text\": \"replace\"}, "
                                "\"add_shapes\": [{\"slide\": 1, \"type\": \"rect\", \"x\": 1, \"y\": 2, \"w\": 3, \"h\": 1, "
                                "\"text\": \"Hello\", \"fill_color\": \"#FF0000\"}], "
                                "\"add_notes\": {\"slide\": 1, \"notes\": \"Notes text\"}, "
                                "\"set_background\": {\"slide\": \"all\", \"color\": \"#1a1a2e\"}, "
                                "\"add_image\": {\"slide\": 1, \"image_path\": \"/path/to/img.png\", "
                                "\"x\": 1, \"y\": 2, \"w\": 4, \"h\": 3}}"
                            )
                        }
                    },
                    "required": ["path", "operations_json"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "pptx_info",
                "description": (
                    "Get detailed metadata and statistics about a PowerPoint file: "
                    "slide count, dimensions, shapes count, text stats, images, tables, "
                    "speaker notes, and file metadata (author, title, dates)."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .pptx file"}
                    },
                    "required": ["path"]
                }
            }
        },

        # ── XLSX TOOLS ──
        {
            "type": "function",
            "function": {
                "name": "read_xlsx",
                "description": (
                    "Read an Excel (.xlsx) file and return content as formatted text tables. "
                    "Shows headers, data rows, cell values with types (numbers, dates, formulas). "
                    "Supports reading specific sheets and limiting rows."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .xlsx file"},
                        "sheet": {"type": "string", "description": "Sheet name to read (default: all sheets)"},
                        "max_rows": {"type": "integer", "description": "Max rows per sheet (default: 100)"},
                        "header": {"type": "boolean", "description": "First row is header (default: true)"}
                    },
                    "required": ["path"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "create_xlsx",
                "description": (
                    "Create an Excel (.xlsx) spreadsheet with data. Supports multiple sheets, "
                    "formatted headers with colors, column widths, auto-filter, freeze panes, "
                    "formulas, number formats, and styled data tables."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Output .xlsx file path"},
                        "title": {"type": "string", "description": "Simple mode title (creates single sheet with title)"},
                        "sheets_json": {
                            "type": "string",
                            "description": (
                                "JSON with sheet definitions:\n"
                                "{\"sheets\": [{"
                                "\"name\": \"Sheet1\", "
                                "\"headers\": [\"Name\", \"Age\"], "
                                "\"data\": [[\"Alice\", 25], [\"Bob\", 30]], "
                                "\"formulas\": {\"C4\": \"=SUM(C2:C3)\"}, "
                                "\"column_widths\": {\"A\": 20, \"B\": 10}, "
                                "\"freeze_panes\": \"A2\", "
                                "\"auto_filter\": true, "
                                "\"bold_headers\": true, "
                                "\"number_format\": {\"C\": \"0.00\"}, "
                                "\"header_color\": \"#4472C4\", "
                                "\"header_text_color\": \"#FFFFFF\"}]}"
                            )
                        }
                    },
                    "required": ["path"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "edit_xlsx",
                "description": (
                    "Edit an existing Excel (.xlsx) file with comprehensive operations: "
                    "add/delete/rename sheets, add/update rows and cells, set formulas, "
                    "add charts (bar/line/pie), format cells, merge cells, sort data, "
                    "find and replace text, set column widths and row heights."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to existing .xlsx file"},
                        "output": {"type": "string", "description": "Output path (default: overwrite)"},
                        "operations_json": {
                            "type": "string",
                            "description": (
                                "JSON object with operations:\n"
                                "{\"add_sheet\": {\"name\": \"New\", \"data\": [[1,2]]}, "
                                "\"delete_sheet\": \"Sheet2\", "
                                "\"rename_sheet\": {\"old\": \"S1\", \"new\": \"S2\"}, "
                                "\"add_rows\": {\"sheet\": \"S1\", \"data\": [[1,2]], \"start_row\": 10}, "
                                "\"update_cell\": {\"sheet\": \"S1\", \"cell\": \"A1\", \"value\": \"new\"}, "
                                "\"add_formula\": {\"sheet\": \"S1\", \"cell\": \"C1\", \"formula\": \"=SUM(A1:B1)\"}, "
                                "\"set_column_width\": {\"sheet\": \"S1\", \"column\": \"A\", \"width\": 25}, "
                                "\"find_replace\": {\"sheet\": \"S1\", \"find\": \"old\", \"replace\": \"new\"}, "
                                "\"sort_range\": {\"sheet\": \"S1\", \"range\": \"A1:C10\", \"by_column\": \"A\", \"ascending\": true}, "
                                "\"merge_cells\": {\"sheet\": \"S1\", \"range\": \"A1:B1\"}, "
                                "\"add_chart\": {\"sheet\": \"S1\", \"chart_type\": \"bar\", \"title\": \"Chart\", "
                                "\"data_range\": \"A1:B10\", \"position\": \"D2\"}, "
                                "\"set_cell_format\": {\"sheet\": \"S1\", \"range\": \"A1:A10\", "
                                "\"bold\": true, \"font_size\": 14, \"bg_color\": \"#FFFF00\", \"alignment\": \"center\"}}"
                            )
                        }
                    },
                    "required": ["path", "operations_json"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "xlsx_info",
                "description": (
                    "Get detailed info about an Excel file: properties, sheets, dimensions, "
                    "data type statistics (numbers, strings, dates, formulas, booleans), "
                    "merged cells, freeze panes, auto-filter, and column widths."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to .xlsx file"}
                    },
                    "required": ["path"]
                }
            }
        },

        # ── DOCX EDIT TOOL ──
        {
            "type": "function",
            "function": {
                "name": "edit_docx",
                "description": (
                    "Edit an existing Word (.docx) file. Can replace text, append content, "
                    "insert headings, add tables, set document properties, delete paragraphs, "
                    "add page breaks, and change font formatting. Multiple operations in one call."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to existing .docx file"},
                        "output": {"type": "string", "description": "Output path (default: overwrite)"},
                        "operations_json": {
                            "type": "string",
                            "description": (
                                "JSON object with operations:\n"
                                "{\"replace_text\": {\"find\": \"old\", \"replace\": \"new\"}, "
                                "\"append_content\": \"# New Section\\nNew paragraph\", "
                                "\"add_table\": {\"headers\": [\"Col1\"], \"rows\": [[\"val\"]]}, "
                                "\"set_properties\": {\"title\": \"My Doc\", \"author\": \"Me\"}, "
                                "\"delete_paragraphs\": {\"index_start\": 5, \"index_end\": 10}, "
                                "\"add_page_break\": {}, "
                                "\"set_font\": {\"paragraph_index\": 0, \"bold\": true, \"size\": 14, \"color\": \"FF0000\"}}"
                            )
                        }
                    },
                    "required": ["path", "operations_json"]
                }
            }
        },

        # ── CSV TOOLS ──
        {
            "type": "function",
            "function": {
                "name": "read_csv",
                "description": (
                    "Read a CSV file and return formatted content with headers, row data, "
                    "and file statistics. Auto-detects encoding (UTF-8, Latin-1, etc.)."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to CSV file"},
                        "delimiter": {"type": "string", "description": "Column delimiter (default: comma)"},
                        "max_rows": {"type": "integer", "description": "Max rows to display (default: 200)"},
                        "encoding": {"type": "string", "description": "File encoding (default: auto-detect)"}
                    },
                    "required": ["path"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "create_csv",
                "description": (
                    "Create a CSV file from structured data."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Output CSV path"},
                        "data_json": {"type": "string", "description": "JSON array of rows: [[\"val1\",\"val2\"],[\"val3\",\"val4\"]]"},
                        "headers_json": {"type": "string", "description": "JSON array of column headers: [\"Col1\",\"Col2\"]"},
                        "delimiter": {"type": "string", "description": "Delimiter character (default: comma)"},
                        "encoding": {"type": "string", "description": "Output encoding (default: utf-8)"}
                    },
                    "required": ["path"]
                }
            }
        },

        # ── CONVERSION TOOL ──
        {
            "type": "function",
            "function": {
                "name": "convert_document",
                "description": (
                    "Convert between document formats. Supports: XLSX->CSV, CSV->XLSX, "
                    "Markdown->DOCX, Markdown->PDF, TXT->DOCX, JSON->XLSX, JSON->CSV. "
                    "All conversions are real and produce actual files."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "input_path": {"type": "string", "description": "Input file path"},
                        "output_path": {"type": "string", "description": "Output file path (determines target format)"},
                        "options": {"type": "string", "description": "JSON string with additional options"}
                    },
                    "required": ["input_path", "output_path"]
                }
            }
        },
    ]


def execute_doc_tool(name: str, args: dict) -> str:
    """Execute a document tool by name."""
    dispatch = {
        'read_pptx': lambda: read_pptx(args['path']),
        'create_pptx': lambda: create_pptx(
            args['path'],
            args.get('title', ''),
            args.get('slides_json', '')
        ),
        'edit_pptx': lambda: edit_pptx(
            args['path'],
            args.get('output', ''),
            args.get('operations_json', '')
        ),
        'pptx_info': lambda: pptx_info(args['path']),

        'read_xlsx': lambda: read_xlsx(
            args['path'],
            args.get('sheet', ''),
            args.get('max_rows', 100),
            args.get('header', True)
        ),
        'create_xlsx': lambda: create_xlsx(
            args['path'],
            args.get('sheets_json', ''),
            args.get('title', '')
        ),
        'edit_xlsx': lambda: edit_xlsx(
            args['path'],
            args.get('output', ''),
            args.get('operations_json', '')
        ),
        'xlsx_info': lambda: xlsx_info(args['path']),

        'edit_docx': lambda: edit_docx(
            args['path'],
            args.get('output', ''),
            args.get('operations_json', '')
        ),

        'read_csv': lambda: read_csv(
            args['path'],
            args.get('delimiter', ','),
            args.get('max_rows', 200),
            args.get('encoding', '')
        ),
        'create_csv': lambda: create_csv(
            args['path'],
            args.get('data_json', ''),
            args.get('headers_json', ''),
            args.get('delimiter', ','),
            args.get('encoding', 'utf-8')
        ),

        'convert_document': lambda: convert_document(
            args['input_path'],
            args['output_path'],
            args.get('options', '')
        ),
    }

    handler = dispatch.get(name)
    if handler:
        return handler()
    return f"[ERROR] Unknown doc tool: {name}"
